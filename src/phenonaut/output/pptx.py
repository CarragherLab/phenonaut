# Copyright © The University of Edinburgh, 2022.
# Development has been supported by GSK.

import datetime
from pathlib import Path
from typing import Optional, Union

import pkg_resources

try:
    import pptx
except ImportError as ie:
    raise (
        f"Import Error - python-pptx functionality called, but it is not installed {ie}"
    )


class PhenonautPPTX:
    """Phenonaut PPTX class

    Allows the creation of PPTX presentation files through insertion of images into slides.

    Instantiate the object using:

    Parameters
    ----------
    cover_title : str, optional
        Title to be displayed on the main title page of the presentation,
        by default "Phenonaut".
    cover_subtitle : str, optional
        Subtitle to be displayed on the main title page of the presentation,
        by default an empty string ("").
    template_file : Optional[Union[Path, str]], optional
        Template file that can be used to style the presentation. Phenonaut
        is supplied with a suitable template which is used as default when
        this argument is None. By default None.
    experiment_hash : Optional[str], optional
        Optionally supply a hash string which is displayed within the PPTX
        file. By default None.


    Then call add_image_slide to add slides, before finally calling save.
    """

    def __init__(
        self,
        cover_title: str = "Phenonaut",
        cover_subtitle: str = "",
        template_file: Optional[Union[Path, str]] = None,
        experiment_hash: Optional[str] = None,
    ):
        """Constructor of PhenonautPPTX object

        Parameters
        ----------
        cover_title : str, optional
            Title to be displayed on the main title page of the presentation,
            by default "Phenonaut".
        cover_subtitle : str, optional
            Subtitle to be displayed on the main title page of the presentation,
            by default an empty string ("").
        template_file : Optional[Union[Path, str]], optional
            Template file that can be used to style the presentation. Phenonaut
            is supplied with a suitable template which is used as default when
            this argument is None. By default None.
        experiment_hash : Optional[str], optional
            Optionally supply a hash string which is displayed within the PPTX
            file. By default None.
        """

        self.top_left_content_loc = (pptx.util.Cm(2.33), pptx.util.Cm(3.31))

        if template_file is None:
            template_file = pkg_resources.resource_filename(
                "phenonaut", "output/ph0_pptx_template.pptx"
            )
        if isinstance(template_file, str):
            template_file = Path(template_file)

        self.pres = pptx.Presentation(template_file)
        slide = self.pres.slides.add_slide(self.pres.slide_layouts[0])
        slide.shapes.title.text = cover_title
        slide.placeholders[1].text = cover_subtitle
        slide.placeholders[
            13
        ].text = f"Generated: {datetime.datetime.now().strftime('%d/%m/%Y, %H:%M:%S')}"
        slide.placeholders[14].text = (
            experiment_hash if experiment_hash is not None else "Hash: NA"
        )

    def add_image_slide(
        self,
        title: str,
        img_file: Union[Path, str],
        width: Optional[Union[int, float]] = None,
        height: Optional[Union[int, float]] = None,
    ):
        """Add an image slide to an existing presentation

        Parameters
        ----------
        title : str
            Slide title
        img_file : Union[Path, str]
            Image file to insert into slide
        width : Optional[Union[int, float]], optional
            Width in cm to scale image. No scaling is done if argument is None,
            by default None.
        height : Optional[Union[int, float]], optional
            Height in cm to scale image. No scaling is done if argument is None,
            by default None.
        """
        if isinstance(img_file, str):
            img_file = Path(img_file)
        if isinstance(width, (int, float)):
            width = pptx.util.Cm(width)
        if isinstance(height, (int, float)):
            height = pptx.util.Cm(height)
        slide = self.pres.slides.add_slide(self.pres.slide_layouts[1])
        slide.placeholders[0].text = title
        slide.shapes.add_picture(
            str(img_file), *self.top_left_content_loc, width=width, height=height
        )

    def save(self, output_pptx_file: Union[Path, str]):
        """Save the presentation

        Parameters
        ----------
        output_pptx_file : Union[Path, str]
            Output file path
        """
        self.output_pptx_file = output_pptx_file
        if isinstance(self.output_pptx_file, str):
            self.output_pptx_file = Path(self.output_pptx_file)
        self.pres.save(self.output_pptx_file)
